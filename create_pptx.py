"""
Generate a detailed product presentation for
Infosys Cobalt Powered AI Contract Lifecycle Management
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Brand colors ──
COBALT_BLUE = RGBColor(0x00, 0x7C, 0xC3)
COBALT_PURPLE = RGBColor(0x5B, 0x2D, 0x8E)
DARK_NAVY = RGBColor(0x0F, 0x17, 0x2A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG = RGBColor(0xF0, 0xF5, 0xFA)
GRAY_TEXT = RGBColor(0x64, 0x74, 0x8B)
DARK_TEXT = RGBColor(0x1E, 0x29, 0x3B)
RED = RGBColor(0xDC, 0x26, 0x26)
GREEN = RGBColor(0x16, 0xA3, 0x4A)
ORANGE = RGBColor(0xEA, 0x58, 0x0C)
AMBER = RGBColor(0xD9, 0x77, 0x06)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
W = prs.slide_width
H = prs.slide_height


def add_gradient_rect(slide, left, top, width, height, color1=COBALT_BLUE, color2=COBALT_PURPLE):
    """Add a colored rectangle as background."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color1
    shape.line.fill.background()
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=18, bold=False,
                 color=DARK_TEXT, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    """Add a text box with formatted text."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_bullet_list(slide, left, top, width, height, items, font_size=14, color=DARK_TEXT, spacing=Pt(6)):
    """Add a bullet list."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_after = spacing
        p.level = 0
    return txBox


def add_card(slide, left, top, width, height, title, body, icon="",
             title_color=COBALT_BLUE, bg_color=WHITE):
    """Add a card-style shape with title and body."""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.color.rgb = RGBColor(0xE2, 0xE8, 0xF0)
    shape.line.width = Pt(1)
    shape.shadow.inherit = False

    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.2)
    tf.margin_right = Inches(0.2)
    tf.margin_top = Inches(0.15)

    p = tf.paragraphs[0]
    p.text = f"{icon} {title}" if icon else title
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = title_color
    p.font.name = "Calibri"
    p.space_after = Pt(6)

    p2 = tf.add_paragraph()
    p2.text = body
    p2.font.size = Pt(11)
    p2.font.color.rgb = GRAY_TEXT
    p2.font.name = "Calibri"
    return shape


def slide_header(slide, title, subtitle=""):
    """Add consistent header bar to a slide."""
    add_gradient_rect(slide, 0, 0, W, Inches(1.1))
    add_text_box(slide, Inches(0.6), Inches(0.15), Inches(10), Inches(0.5),
                 title, font_size=28, bold=True, color=WHITE)
    if subtitle:
        add_text_box(slide, Inches(0.6), Inches(0.6), Inches(10), Inches(0.4),
                     subtitle, font_size=14, color=RGBColor(0xBF, 0xDB, 0xFE))
    # Infosys Cobalt badge
    add_text_box(slide, Inches(10.5), Inches(0.25), Inches(2.5), Inches(0.5),
                 "Infosys Cobalt", font_size=14, bold=True, color=WHITE, alignment=PP_ALIGN.RIGHT)


# =====================================================================
# SLIDE 1 — Title Slide
# =====================================================================
slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
add_gradient_rect(slide1, 0, 0, W, H, DARK_NAVY)

# Accent line
shape = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(2.8), Inches(1), Inches(0.06))
shape.fill.solid()
shape.fill.fore_color.rgb = COBALT_BLUE
shape.line.fill.background()

add_text_box(slide1, Inches(0.8), Inches(1.2), Inches(11), Inches(0.6),
             "Infosys Cobalt", font_size=20, color=RGBColor(0x60, 0xA5, 0xFA), bold=True)
add_text_box(slide1, Inches(0.8), Inches(1.8), Inches(11), Inches(1),
             "AI Contract Lifecycle Management", font_size=44, bold=True, color=WHITE)
add_text_box(slide1, Inches(0.8), Inches(3.0), Inches(10), Inches(0.8),
             "Improve productivity of document reviews by identifying key elements,\n"
             "contract draft generation, risk management, and compliance automation.",
             font_size=16, color=RGBColor(0x94, 0xA3, 0xB8))
add_text_box(slide1, Inches(0.8), Inches(4.2), Inches(5), Inches(0.4),
             "Powered by OpenAI GPT-4o  |  Streamlit  |  Cloud-Native",
             font_size=13, color=RGBColor(0x60, 0xA5, 0xFA))

# KPI highlights at bottom
for i, (val, label) in enumerate([("4", "AI Agents"), ("12", "Feature Pages"),
                                    ("6", "Contract Templates"), ("20+", "Clause Categories")]):
    x = Inches(0.8 + i * 2.8)
    add_text_box(slide1, x, Inches(5.4), Inches(2), Inches(0.5),
                 val, font_size=36, bold=True, color=COBALT_BLUE)
    add_text_box(slide1, x, Inches(5.95), Inches(2), Inches(0.3),
                 label, font_size=13, color=RGBColor(0x94, 0xA3, 0xB8))


# =====================================================================
# SLIDE 2 — Problem Statement & Value Proposition
# =====================================================================
slide2 = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(slide2, "The Challenge", "Why enterprises need AI-powered contract management")

# Problem cards (left)
add_text_box(slide2, Inches(0.6), Inches(1.3), Inches(5), Inches(0.4),
             "Current Pain Points", font_size=18, bold=True, color=RED)
problems = [
    ("Manual Review", "Contracts take 2-4 weeks for manual review, delaying deals and increasing costs."),
    ("Hidden Risks", "Critical risky clauses and compliance gaps go undetected until disputes arise."),
    ("Inconsistent Drafts", "Each attorney drafts differently, leading to inconsistent terms and missing protections."),
    ("No Visibility", "No centralized view of contract portfolio, expiration dates, or risk exposure."),
]
for i, (title, body) in enumerate(problems):
    add_card(slide2, Inches(0.6), Inches(1.9 + i * 1.25), Inches(5.5), Inches(1.05),
             title, body, icon="⚠️", title_color=RED, bg_color=RGBColor(0xFE, 0xF2, 0xF2))

# Solution cards (right)
add_text_box(slide2, Inches(7), Inches(1.3), Inches(5), Inches(0.4),
             "Our Solution", font_size=18, bold=True, color=GREEN)
solutions = [
    ("AI Extraction", "ClauseScout extracts all key elements in seconds — parties, dates, obligations, penalties."),
    ("Risk Scoring", "RiskRadar scores every contract 0-100 with detailed clause-level risk breakdown."),
    ("Smart Drafting", "DraftCraft generates legally sound contracts from templates + clause library."),
    ("Full Visibility", "Real-time dashboard with KPIs, expiration alerts, audit trail, and email notifications."),
]
for i, (title, body) in enumerate(solutions):
    add_card(slide2, Inches(7), Inches(1.9 + i * 1.25), Inches(5.5), Inches(1.05),
             title, body, icon="✅", title_color=GREEN, bg_color=RGBColor(0xF0, 0xFD, 0xF4))


# =====================================================================
# SLIDE 3 — AI Agents Overview
# =====================================================================
slide3 = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(slide3, "4 Intelligent AI Agents", "Purpose-built agents powered by OpenAI GPT-4o")

agents = [
    ("🔍 ClauseScout", "Contract Element Extractor",
     "Scans uploaded contracts (PDF, DOCX, scanned images via OCR) and extracts every key element — "
     "parties, dates, obligations, penalties, payment terms, governing law, IP rights, and more. "
     "Generates structured JSON output for downstream processing.",
     ["PDF / DOCX / TXT parsing", "Tesseract OCR for scanned docs", "Structured JSON extraction", "Quick summary generation"]),
    ("✍️ DraftCraft", "Contract Draft Generator",
     "Generates professionally structured, legally sound contract documents from 6 built-in templates "
     "with AI-powered customization. Supports iterative refinement and clause library integration.",
     ["6 templates (NDA, MSA, SOW, etc.)", "Custom clause insertion", "Iterative AI refinement", "Export to PDF / DOCX"]),
    ("🛡️ RiskRadar", "Risk Analysis Specialist",
     "Evaluates contracts for risky clauses, compliance issues, and missing protections. "
     "Scores contracts 0-100 and provides actionable negotiation recommendations.",
     ["0-100 risk scoring", "Clause-level risk breakdown", "Compliance checks (GDPR, CCPA)", "Negotiation recommendations"]),
    ("⚖️ DiffLens", "Contract Comparison Analyst",
     "Performs side-by-side AI analysis of two contracts, identifying key differences, "
     "added/removed/modified clauses, and recommending which terms are more favorable.",
     ["Side-by-side comparison", "Structural text diff", "AI-powered analysis", "Favorability assessment"]),
]

for i, (name, role, desc, features) in enumerate(agents):
    x = Inches(0.5 + i * 3.15)
    # Card background
    card = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(1.3), Inches(2.95), Inches(5.7))
    card.fill.solid()
    card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = RGBColor(0xE2, 0xE8, 0xF0)
    card.line.width = Pt(1)

    # Agent icon header
    header = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(1.3), Inches(2.95), Inches(0.9))
    header.fill.solid()
    header.fill.fore_color.rgb = COBALT_BLUE if i % 2 == 0 else COBALT_PURPLE
    header.line.fill.background()
    tf = header.text_frame
    tf.margin_left = Inches(0.15)
    tf.margin_top = Inches(0.1)
    p = tf.paragraphs[0]
    p.text = name
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = "Calibri"
    p2 = tf.add_paragraph()
    p2.text = role
    p2.font.size = Pt(10)
    p2.font.color.rgb = RGBColor(0xBF, 0xDB, 0xFE)
    p2.font.name = "Calibri"

    # Description
    add_text_box(slide3, x + Inches(0.15), Inches(2.35), Inches(2.65), Inches(1.8),
                 desc, font_size=10, color=GRAY_TEXT)

    # Features
    add_bullet_list(slide3, x + Inches(0.15), Inches(4.2), Inches(2.65), Inches(2.5),
                    features, font_size=10, color=DARK_TEXT, spacing=Pt(4))


# =====================================================================
# SLIDE 4 — Platform Architecture
# =====================================================================
slide4 = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(slide4, "Platform Architecture", "Cloud-native, containerized, enterprise-ready")

# Architecture layers
layers = [
    ("Presentation Layer", "Streamlit Web Application", COBALT_BLUE,
     ["Responsive UI with Infosys Cobalt branding", "Role-based access (Admin / Analyst / Viewer)",
      "Interactive dashboards with Plotly charts", "Multi-format export (PDF, DOCX, Excel)"]),
    ("AI / Intelligence Layer", "OpenAI GPT-4o", COBALT_PURPLE,
     ["4 specialized AI agents", "Structured JSON extraction with response_format",
      "Iterative refinement with conversation history", "OCR pipeline (Tesseract + pdf2image)"]),
    ("Business Logic Layer", "Python / FastAPI", RGBColor(0x05, 0x96, 0x69),
     ["Contract lifecycle state machine", "Clause library with usage tracking",
      "Email alert engine (SMTP/TLS)", "Comprehensive audit trail"]),
    ("Data & Storage Layer", "SQLite + File System", RGBColor(0xD9, 0x77, 0x06),
     ["SQLite with WAL mode for concurrent access", "5 tables: contracts, drafts, risk_analyses, clauses, audit",
      "File uploads with OCR-enabled parsing", "bcrypt-hashed authentication credentials"]),
]

for i, (title, tech, color, features) in enumerate(layers):
    y = Inches(1.4 + i * 1.45)
    # Layer bar
    bar = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), y, Inches(12.3), Inches(1.3))
    bar.fill.solid()
    bar.fill.fore_color.rgb = WHITE
    bar.line.color.rgb = color
    bar.line.width = Pt(2)

    # Color tag
    tag = slide4.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), y, Inches(0.12), Inches(1.3))
    tag.fill.solid()
    tag.fill.fore_color.rgb = color
    tag.line.fill.background()

    add_text_box(slide4, Inches(0.9), y + Inches(0.08), Inches(3), Inches(0.3),
                 title, font_size=14, bold=True, color=color)
    add_text_box(slide4, Inches(0.9), y + Inches(0.38), Inches(3), Inches(0.25),
                 tech, font_size=10, color=GRAY_TEXT)

    for j, feat in enumerate(features):
        col = j % 2
        row = j // 2
        fx = Inches(4.2 + col * 4.3)
        fy = y + Inches(0.1 + row * 0.55)
        add_text_box(slide4, fx, fy, Inches(4), Inches(0.45),
                     f"▸  {feat}", font_size=10, color=DARK_TEXT)


# =====================================================================
# SLIDE 5 — Feature Deep Dive: Upload & Review + OCR
# =====================================================================
slide5 = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(slide5, "Upload & Review with OCR", "ClauseScout Agent — intelligent document processing")

# Flow steps
steps = [
    ("1", "Upload", "PDF, DOCX, TXT, or\nscanned image"),
    ("2", "Smart Parse", "Standard text extraction\nwith OCR fallback"),
    ("3", "AI Extract", "GPT-4o identifies all\nkey contract elements"),
    ("4", "Review", "Structured display of\nparties, dates, terms"),
    ("5", "Save", "Store to repository\nwith metadata & tags"),
]

for i, (num, title, desc) in enumerate(steps):
    x = Inches(0.5 + i * 2.5)
    # Circle number
    circle = slide5.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.7), Inches(1.5), Inches(0.6), Inches(0.6))
    circle.fill.solid()
    circle.fill.fore_color.rgb = COBALT_BLUE
    circle.line.fill.background()
    tf = circle.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Inches(0)
    p = tf.paragraphs[0]
    p.text = num
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # Arrow (except last)
    if i < len(steps) - 1:
        arrow = slide5.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x + Inches(1.6), Inches(1.62), Inches(0.6), Inches(0.35))
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = RGBColor(0xCB, 0xD5, 0xE1)
        arrow.line.fill.background()

    add_text_box(slide5, x, Inches(2.2), Inches(2.1), Inches(0.3),
                 title, font_size=14, bold=True, color=COBALT_BLUE, alignment=PP_ALIGN.CENTER)
    add_text_box(slide5, x, Inches(2.55), Inches(2.1), Inches(0.6),
                 desc, font_size=10, color=GRAY_TEXT, alignment=PP_ALIGN.CENTER)

# Extracted elements showcase
add_text_box(slide5, Inches(0.6), Inches(3.5), Inches(5), Inches(0.4),
             "Extracted Key Elements", font_size=16, bold=True, color=DARK_TEXT)

elements = [
    ("Parties", "Names, roles (buyer/seller/landlord)"),
    ("Dates", "Effective, expiration, renewal terms"),
    ("Obligations", "Per-party contractual duties"),
    ("Penalties", "Breach penalties, liquidated damages"),
    ("Payment Terms", "Amounts, schedules, currency"),
    ("Governing Law", "Jurisdiction, dispute resolution"),
    ("IP Rights", "Ownership, licenses, assignments"),
    ("Confidentiality", "NDA terms, exclusions, duration"),
    ("Force Majeure", "Excusable events, notice periods"),
    ("Termination", "Notice periods, cure rights, severance"),
]

for i, (elem, desc) in enumerate(elements):
    col = i % 3
    row = i // 3
    x = Inches(0.6 + col * 4.2)
    y = Inches(4.0 + row * 0.85)
    add_card(slide5, x, y, Inches(3.9), Inches(0.7), elem, desc, title_color=COBALT_BLUE)


# =====================================================================
# SLIDE 6 — Risk Analysis Deep Dive
# =====================================================================
slide6 = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(slide6, "AI Risk Analysis", "RiskRadar Agent — comprehensive contract risk assessment")

# Risk scoring scale
add_text_box(slide6, Inches(0.6), Inches(1.3), Inches(5), Inches(0.4),
             "Risk Scoring Scale (0-100)", font_size=16, bold=True, color=DARK_TEXT)

risk_levels = [
    ("Low (1-20)", "Well-balanced contract", GREEN, Inches(0.6)),
    ("Medium (21-40)", "Minor concerns", AMBER, Inches(3.3)),
    ("High (41-70)", "Significant risks", ORANGE, Inches(6.0)),
    ("Critical (71-100)", "Immediate action needed", RED, Inches(8.7)),
]

for label, desc, color, x in risk_levels:
    bar = slide6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(1.8), Inches(2.4), Inches(0.8))
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()
    tf = bar.text_frame
    tf.margin_left = Inches(0.15)
    tf.margin_top = Inches(0.08)
    p = tf.paragraphs[0]
    p.text = label
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p2 = tf.add_paragraph()
    p2.text = desc
    p2.font.size = Pt(9)
    p2.font.color.rgb = WHITE

# Analysis output sections
add_text_box(slide6, Inches(0.6), Inches(2.9), Inches(12), Inches(0.4),
             "What RiskRadar Analyzes", font_size=16, bold=True, color=DARK_TEXT)

analysis_items = [
    ("⚠️ Risky Clauses", "Identifies unfavorable terms like\nunlimited liability, auto-renewal traps,\none-sided indemnification"),
    ("📋 Compliance Issues", "Checks against GDPR, CCPA, SOX,\nindustry regulations, and\njurisdiction conflicts"),
    ("🚫 Missing Protections", "Flags absent standard clauses:\nforce majeure, IP protection,\nlimitation of liability, data protection"),
    ("✅ Favorable Clauses", "Highlights terms that benefit\nthe reviewing party, such as\ncaps, carve-outs, and protections"),
    ("💬 Negotiation Points", "Provides prioritized suggestions\nfor contract modifications with\nrecommended alternative language"),
    ("📊 Executive Summary", "3-5 sentence overview suitable\nfor leadership briefings and\nquick decision-making"),
]

for i, (title, body) in enumerate(analysis_items):
    col = i % 3
    row = i // 3
    x = Inches(0.5 + col * 4.2)
    y = Inches(3.4 + row * 2.0)
    add_card(slide6, x, y, Inches(3.9), Inches(1.7), title, body, title_color=COBALT_BLUE)


# =====================================================================
# SLIDE 7 — Draft Generation & Clause Library
# =====================================================================
slide7 = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(slide7, "Draft Generation & Clause Library", "DraftCraft Agent + reusable clause management")

# Templates
add_text_box(slide7, Inches(0.6), Inches(1.3), Inches(5), Inches(0.4),
             "6 Built-in Contract Templates", font_size=16, bold=True, color=DARK_TEXT)

templates = [
    ("NDA", "Non-Disclosure\nAgreement"),
    ("MSA", "Master Service\nAgreement"),
    ("SOW", "Statement\nof Work"),
    ("Employment", "Employment\nAgreement"),
    ("Vendor", "Vendor\nAgreement"),
    ("Lease", "Lease\nAgreement"),
]

for i, (short, full) in enumerate(templates):
    x = Inches(0.5 + i * 2.1)
    card = slide7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(1.8), Inches(1.9), Inches(1.1))
    card.fill.solid()
    card.fill.fore_color.rgb = COBALT_BLUE if i % 2 == 0 else COBALT_PURPLE
    card.line.fill.background()
    tf = card.text_frame
    tf.margin_left = Inches(0.15)
    tf.margin_top = Inches(0.1)
    p = tf.paragraphs[0]
    p.text = short
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.text = full
    p2.font.size = Pt(9)
    p2.font.color.rgb = RGBColor(0xBF, 0xDB, 0xFE)
    p2.alignment = PP_ALIGN.CENTER

# Draft workflow
add_text_box(slide7, Inches(0.6), Inches(3.2), Inches(6), Inches(0.4),
             "AI-Powered Drafting Workflow", font_size=16, bold=True, color=DARK_TEXT)

workflow = [
    "Select contract type and fill in parameters (parties, dates, jurisdiction)",
    "Choose reusable clauses from the Clause Library to include",
    "DraftCraft generates a complete, legally structured draft",
    "Review, edit inline, and request AI-powered refinements",
    "Export as PDF or DOCX, or save to the contract repository",
]
add_bullet_list(slide7, Inches(0.6), Inches(3.7), Inches(6), Inches(2.5),
                [f"{i+1}.  {w}" for i, w in enumerate(workflow)], font_size=12, spacing=Pt(8))

# Clause Library
add_text_box(slide7, Inches(7.2), Inches(3.2), Inches(5), Inches(0.4),
             "📚 Clause Library (20 Categories)", font_size=16, bold=True, color=DARK_TEXT)

categories = [
    "Termination", "Indemnification", "Limitation of Liability",
    "Confidentiality", "Force Majeure", "Intellectual Property",
    "Non-Compete", "Data Protection / GDPR", "Payment Terms",
    "Warranty", "Dispute Resolution", "Governing Law",
    "Insurance", "Compliance", "Severability",
]
add_bullet_list(slide7, Inches(7.2), Inches(3.7), Inches(5.5), Inches(3.5),
                categories, font_size=11, color=DARK_TEXT, spacing=Pt(3))


# =====================================================================
# SLIDE 8 — Enterprise Features
# =====================================================================
slide8 = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(slide8, "Enterprise Features", "Authentication, audit trail, email alerts, and bulk processing")

features = [
    ("👥 Multi-User Authentication", "Role-Based Access Control",
     ["3 roles: Admin, Analyst, Viewer",
      "bcrypt-hashed passwords with cookie sessions",
      "Admin can add/remove users and manage roles",
      "Viewers can browse; Analysts can upload and analyze; Admins have full control"]),
    ("📋 Audit Trail", "Complete Activity Tracking",
     ["Tracks 16 action types across the system",
      "Who uploaded, analyzed, drafted, compared, deleted",
      "Per-contract activity history",
      "Exportable audit logs for compliance reporting"]),
    ("📧 Email Alerts", "Proactive Notifications",
     ["SMTP/TLS email delivery",
      "Branded HTML emails with Infosys Cobalt styling",
      "Configurable alert window (7-90 days before expiry)",
      "High-risk contract alerts with score and top risks"]),
    ("📦 Bulk Upload", "Batch Processing at Scale",
     ["Upload multiple files simultaneously",
      "Automatic AI extraction per file",
      "Batch risk scoring across all uploaded contracts",
      "OCR fallback for scanned documents in batch"]),
]

for i, (title, subtitle, items) in enumerate(features):
    col = i % 2
    row = i // 2
    x = Inches(0.5 + col * 6.3)
    y = Inches(1.4 + row * 3.0)

    card = slide8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(6.0), Inches(2.7))
    card.fill.solid()
    card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = COBALT_BLUE
    card.line.width = Pt(1.5)

    add_text_box(slide8, x + Inches(0.2), y + Inches(0.1), Inches(5.5), Inches(0.35),
                 title, font_size=15, bold=True, color=COBALT_BLUE)
    add_text_box(slide8, x + Inches(0.2), y + Inches(0.45), Inches(5.5), Inches(0.25),
                 subtitle, font_size=11, color=GRAY_TEXT)
    add_bullet_list(slide8, x + Inches(0.3), y + Inches(0.75), Inches(5.4), Inches(1.8),
                    [f"✓  {item}" for item in items], font_size=11, spacing=Pt(3))


# =====================================================================
# SLIDE 9 — Deployment Options
# =====================================================================
slide9 = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(slide9, "Cloud Deployment", "Multiple deployment options for every environment")

deploy_options = [
    ("Streamlit Cloud", "Fastest path to production",
     ["One-click deploy from GitHub", "Auto-redeploy on git push",
      "Built-in secrets management", "Free tier available",
      "packages.txt for system deps (Tesseract)"]),
    ("Docker / Docker Compose", "Containerized deployment",
     ["Python 3.11-slim base image", "Tesseract + Poppler pre-installed",
      "Persistent volumes for data/uploads/config",
      "Health checks and auto-restart",
      "Single command: docker-compose up"]),
    ("AWS / Cloud Infrastructure", "Enterprise-scale deployment",
     ["ECS/Fargate or EKS for orchestration",
      "RDS for database (replace SQLite)",
      "S3 for contract file storage",
      "ALB + Route 53 for DNS",
      "CloudWatch for monitoring"]),
]

for i, (title, subtitle, items) in enumerate(deploy_options):
    x = Inches(0.5 + i * 4.2)
    card = slide9.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(1.4), Inches(3.9), Inches(5.5))
    card.fill.solid()
    card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = COBALT_BLUE
    card.line.width = Pt(1.5)

    # Header
    hdr = slide9.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(1.4), Inches(3.9), Inches(1.0))
    hdr.fill.solid()
    hdr.fill.fore_color.rgb = [COBALT_BLUE, COBALT_PURPLE, RGBColor(0x05, 0x96, 0x69)][i]
    hdr.line.fill.background()
    tf = hdr.text_frame
    tf.margin_left = Inches(0.2)
    tf.margin_top = Inches(0.12)
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p2 = tf.add_paragraph()
    p2.text = subtitle
    p2.font.size = Pt(10)
    p2.font.color.rgb = RGBColor(0xBF, 0xDB, 0xFE)

    add_bullet_list(slide9, x + Inches(0.2), Inches(2.6), Inches(3.5), Inches(4.0),
                    [f"▸  {item}" for item in items], font_size=11, spacing=Pt(8))


# =====================================================================
# SLIDE 10 — Tech Stack
# =====================================================================
slide10 = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(slide10, "Technology Stack", "Modern, proven technologies for reliability and scale")

stack = [
    ("Frontend", [("Streamlit", "Interactive web framework"), ("Plotly", "Rich interactive charts"), ("HTML/CSS", "Custom Cobalt branding")]),
    ("AI Engine", [("OpenAI GPT-4o", "Primary LLM for all agents"), ("GPT-4o-mini", "Fast model for summaries"), ("JSON mode", "Structured AI outputs")]),
    ("OCR", [("Tesseract", "Open-source OCR engine"), ("pdf2image", "PDF to image conversion"), ("Pillow", "Image processing")]),
    ("Backend", [("Python 3.11", "Core runtime"), ("SQLite + WAL", "Embedded database"), ("bcrypt", "Password hashing")]),
    ("Export", [("fpdf2", "PDF generation"), ("python-docx", "DOCX generation"), ("openpyxl", "Excel export")]),
    ("DevOps", [("Docker", "Containerization"), ("GitHub", "Source control + CI"), ("Streamlit Cloud", "Managed hosting")]),
]

for i, (category, techs) in enumerate(stack):
    col = i % 3
    row = i // 3
    x = Inches(0.5 + col * 4.2)
    y = Inches(1.4 + row * 3.0)

    add_text_box(slide10, x, y, Inches(3.8), Inches(0.35),
                 category, font_size=15, bold=True, color=COBALT_BLUE)

    for j, (tech, desc) in enumerate(techs):
        ty = y + Inches(0.45 + j * 0.7)
        card = slide10.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, ty, Inches(3.8), Inches(0.6))
        card.fill.solid()
        card.fill.fore_color.rgb = LIGHT_BG
        card.line.color.rgb = RGBColor(0xE2, 0xE8, 0xF0)
        card.line.width = Pt(0.5)
        tf = card.text_frame
        tf.margin_left = Inches(0.15)
        tf.margin_top = Inches(0.05)
        p = tf.paragraphs[0]
        p.text = tech
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = DARK_TEXT
        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(9)
        p2.font.color.rgb = GRAY_TEXT


# =====================================================================
# SLIDE 11 — ROI & Business Impact
# =====================================================================
slide11 = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(slide11, "Business Impact & ROI", "Measurable value across the contract lifecycle")

metrics = [
    ("90%", "Faster\nContract Review", "From weeks to minutes\nwith AI extraction"),
    ("75%", "Risk Detection\nImprovement", "AI catches risks that\nmanual review misses"),
    ("60%", "Faster Draft\nGeneration", "Templates + AI = instant\nprofessional contracts"),
    ("100%", "Audit\nCompliance", "Every action tracked\nand exportable"),
]

for i, (pct, label, desc) in enumerate(metrics):
    x = Inches(0.5 + i * 3.2)
    card = slide11.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(1.5), Inches(2.9), Inches(2.5))
    card.fill.solid()
    card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = COBALT_BLUE
    card.line.width = Pt(2)

    add_text_box(slide11, x, Inches(1.7), Inches(2.9), Inches(0.7),
                 pct, font_size=44, bold=True, color=COBALT_BLUE, alignment=PP_ALIGN.CENTER)
    add_text_box(slide11, x, Inches(2.4), Inches(2.9), Inches(0.5),
                 label, font_size=13, bold=True, color=DARK_TEXT, alignment=PP_ALIGN.CENTER)
    add_text_box(slide11, x, Inches(3.0), Inches(2.9), Inches(0.7),
                 desc, font_size=10, color=GRAY_TEXT, alignment=PP_ALIGN.CENTER)

# Value drivers
add_text_box(slide11, Inches(0.6), Inches(4.3), Inches(12), Inches(0.4),
             "Key Value Drivers", font_size=16, bold=True, color=DARK_TEXT)

drivers = [
    ("Cost Reduction", "Eliminate manual review bottlenecks.\nReduce outside counsel dependency\nfor routine contract work."),
    ("Risk Mitigation", "Proactively identify risky clauses\nbefore signing. Catch compliance\ngaps that lead to penalties."),
    ("Operational Efficiency", "Centralized repository with search.\nAutomated expiry alerts prevent\nmissed renewal deadlines."),
    ("Decision Speed", "Real-time dashboards give leadership\ninstant visibility into portfolio risk,\nstatus, and contract value."),
]

for i, (title, body) in enumerate(drivers):
    x = Inches(0.5 + i * 3.2)
    add_card(slide11, x, Inches(4.8), Inches(2.9), Inches(2.2), title, body, title_color=COBALT_BLUE)


# =====================================================================
# SLIDE 12 — Roadmap
# =====================================================================
slide12 = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(slide12, "Product Roadmap", "Continuous innovation and enterprise hardening")

phases = [
    ("Phase 1 — Current", "Q1 2025", COBALT_BLUE,
     ["4 AI agents (ClauseScout, DraftCraft, RiskRadar, DiffLens)",
      "6 contract templates + Clause Library",
      "Multi-user auth + Audit Trail",
      "OCR + Bulk Upload",
      "Email alerts + Dashboard"]),
    ("Phase 2 — Next", "Q2 2025", COBALT_PURPLE,
     ["Microsoft 365 / SharePoint integration",
      "Advanced analytics with trend reporting",
      "Multi-language contract support",
      "API endpoints for system integration",
      "Custom template builder"]),
    ("Phase 3 — Future", "Q3-Q4 2025", RGBColor(0x05, 0x96, 0x69),
     ["Enterprise SSO (SAML / OIDC)",
      "PostgreSQL / cloud database backend",
      "S3 / Azure Blob for file storage",
      "Workflow automation (approval chains)",
      "AI-powered contract negotiation assistant"]),
]

for i, (title, timeline, color, items) in enumerate(phases):
    x = Inches(0.5 + i * 4.2)
    card = slide12.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(1.4), Inches(3.9), Inches(5.5))
    card.fill.solid()
    card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = color
    card.line.width = Pt(2)

    hdr = slide12.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(1.4), Inches(3.9), Inches(0.9))
    hdr.fill.solid()
    hdr.fill.fore_color.rgb = color
    hdr.line.fill.background()
    tf = hdr.text_frame
    tf.margin_left = Inches(0.2)
    tf.margin_top = Inches(0.1)
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p2 = tf.add_paragraph()
    p2.text = timeline
    p2.font.size = Pt(11)
    p2.font.color.rgb = RGBColor(0xBF, 0xDB, 0xFE)

    add_bullet_list(slide12, x + Inches(0.2), Inches(2.5), Inches(3.5), Inches(4.2),
                    [f"▸  {item}" for item in items], font_size=12, spacing=Pt(10))


# =====================================================================
# SLIDE 13 — Thank You / Contact
# =====================================================================
slide13 = prs.slides.add_slide(prs.slide_layouts[6])
add_gradient_rect(slide13, 0, 0, W, H, DARK_NAVY)

shape = slide13.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4), Inches(2.5), Inches(1.5), Inches(0.06))
shape.fill.solid()
shape.fill.fore_color.rgb = COBALT_BLUE
shape.line.fill.background()

add_text_box(slide13, Inches(0.8), Inches(1.2), Inches(11), Inches(0.5),
             "Infosys Cobalt", font_size=18, color=RGBColor(0x60, 0xA5, 0xFA), bold=True)
add_text_box(slide13, Inches(0.8), Inches(1.7), Inches(11), Inches(0.8),
             "Thank You", font_size=48, bold=True, color=WHITE)
add_text_box(slide13, Inches(0.8), Inches(2.7), Inches(10), Inches(0.6),
             "AI Contract Lifecycle Management — Transforming how enterprises\n"
             "manage, analyze, and optimize their commercial contracts.",
             font_size=16, color=RGBColor(0x94, 0xA3, 0xB8))

add_text_box(slide13, Inches(0.8), Inches(4.0), Inches(4), Inches(0.4),
             "Live Demo", font_size=14, bold=True, color=RGBColor(0x60, 0xA5, 0xFA))
add_text_box(slide13, Inches(0.8), Inches(4.4), Inches(6), Inches(0.3),
             "contractmgmt.streamlit.app", font_size=14, color=WHITE)

add_text_box(slide13, Inches(0.8), Inches(5.0), Inches(4), Inches(0.4),
             "Source Code", font_size=14, bold=True, color=RGBColor(0x60, 0xA5, 0xFA))
add_text_box(slide13, Inches(0.8), Inches(5.4), Inches(6), Inches(0.3),
             "github.com/ajittgosavii/contractmgmt", font_size=14, color=WHITE)

add_text_box(slide13, Inches(0.8), Inches(6.2), Inches(10), Inches(0.3),
             "Powered by OpenAI GPT-4o  |  Streamlit  |  Infosys Cobalt Cloud Platform",
             font_size=12, color=RGBColor(0x64, 0x74, 0x8B))


# ── Save ──
output_path = os.path.join(os.path.dirname(__file__), "Infosys_Cobalt_AI_CLM_Presentation.pptx")
prs.save(output_path)
print(f"Presentation saved to: {output_path}")
print(f"Total slides: {len(prs.slides)}")
